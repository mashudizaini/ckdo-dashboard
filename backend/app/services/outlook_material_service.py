"""
Outlook Material Service — PAC module.
Stores and retrieves the reference source files (economic reports, market
data, etc.) uploaded ahead of writing the Business Plan Outlook.

Convert stage: each file is summarized into a structured Markdown "brief"
once (convert_material), and that brief — not the raw file — is what
generate_outlook reads from on every generation. This avoids re-reading
potentially large PDFs on every regenerate (slow, expensive, and prone to
drift between runs) in favor of a compact, reusable, point-form reference.
"""
import asyncio
import os
from datetime import datetime
from typing import Optional
from sqlalchemy import select, delete
from sqlalchemy.ext.asyncio import AsyncSession
from app.models.outlook_material import OutlookMaterial
from app.services.ai_service import AIService
import structlog

logger = structlog.get_logger()

_UPLOAD_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.dirname(__file__))), "uploads", "outlook_materials"
)
os.makedirs(_UPLOAD_DIR, exist_ok=True)

# Conservative cap so extracted text + prompt + output stay comfortably
# within the num_ctx window used for the summarization call.
_MAX_EXTRACT_CHARS = 15000


class OutlookMaterialService:

    def storage_path(self, filename: str) -> str:
        return os.path.join(_UPLOAD_DIR, filename)

    async def save_files(
        self,
        db: AsyncSession,
        plan_year: int,
        files: list,
        username: str,
        category: str = "material",
    ) -> dict:
        saved = []
        for file in files:
            content = await file.read()
            ext = os.path.splitext(file.filename or "")[1]
            stored_name = f"{category}_{plan_year}_{datetime.utcnow().strftime('%Y%m%d%H%M%S%f')}{ext}"
            with open(self.storage_path(stored_name), "wb") as f:
                f.write(content)

            row = OutlookMaterial(
                plan_year=plan_year,
                category=category,
                filename=stored_name,
                original_name=file.filename or stored_name,
                content_type=file.content_type,
                file_size=len(content),
                uploaded_by=username,
            )
            db.add(row)
            await db.flush()
            await db.refresh(row)
            saved.append(self._to_dict(row))
        return {"success": True, "count": len(saved), "data": saved}

    async def list_materials(self, db: AsyncSession, plan_year: Optional[int] = None, category: Optional[str] = None) -> dict:
        q = select(OutlookMaterial).order_by(OutlookMaterial.created_at.desc())
        if plan_year:
            q = q.where(OutlookMaterial.plan_year == plan_year)
        if category:
            q = q.where(OutlookMaterial.category == category)
        result = await db.execute(q)
        rows = result.scalars().all()
        return {"success": True, "count": len(rows), "data": [self._to_dict(r) for r in rows]}

    async def get_material(self, db: AsyncSession, material_id: int) -> Optional[OutlookMaterial]:
        return await db.get(OutlookMaterial, material_id)

    async def delete_material(self, db: AsyncSession, material_id: int) -> dict:
        row = await db.get(OutlookMaterial, material_id)
        if not row:
            return {"success": False, "error": "Not found"}
        path = self.storage_path(row.filename)
        if os.path.exists(path):
            os.remove(path)
        await db.execute(delete(OutlookMaterial).where(OutlookMaterial.id == material_id))
        return {"success": True, "message": f"Deleted material #{material_id}"}

    # ── Convert stage ────────────────────────────────────────────────────

    def _extract_text(self, path: str, original_name: str) -> str:
        ext = os.path.splitext(original_name or path)[1].lower()

        if ext == ".pdf":
            import fitz
            doc = fitz.open(path)
            try:
                return "\n".join(page.get_text() for page in doc)
            finally:
                doc.close()

        if ext in (".docx", ".doc"):
            import docx
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)

        if ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        slide_lines.append(shape.text_frame.text)
                    elif shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text for c in row.cells if c.text]
                            if cells:
                                slide_lines.append(" | ".join(cells))
                if slide_lines:
                    lines.append(f"# Slide {i}")
                    lines.extend(slide_lines)
            return "\n".join(lines)

        if ext in (".xlsx", ".xlsm"):
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None]
                    if vals:
                        lines.append(" | ".join(vals))
            return "\n".join(lines)

        if ext in (".txt", ".csv", ".md"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        raise ValueError(f"Format {ext or '(tanpa ekstensi)'} belum didukung untuk convert otomatis")

    async def convert_material(self, db: AsyncSession, material_id: int, provider: str = "onprem", gemini_api_key: str = None) -> dict:
        """Extract text from the uploaded file and summarize it into a
        structured Markdown brief via AI — run once per file, reused on
        every Outlook generation afterwards."""
        row = await db.get(OutlookMaterial, material_id)
        if not row:
            return {"success": False, "error": "Not found"}

        row.brief_status = "converting"
        await db.flush()

        try:
            path = self.storage_path(row.filename)
            if not os.path.exists(path):
                raise FileNotFoundError("File tidak ditemukan di server")

            raw_text = await asyncio.to_thread(self._extract_text, path, row.original_name)
            raw_text = (raw_text or "").strip()
            if not raw_text:
                raise ValueError("Tidak ada teks yang bisa diekstrak dari file ini (kemungkinan hasil scan/gambar)")

            truncated = len(raw_text) > _MAX_EXTRACT_CHARS
            text_for_ai = raw_text[:_MAX_EXTRACT_CHARS]
            truncation_note = "\n[...dipotong, dokumen aslinya lebih panjang dari ini...]" if truncated else ""

            purpose = (
                "acuan STRUKTUR/FORMAT laporan Business Plan Outlook (Global Economic Outlook, "
                "Indonesia Economic Outlook, Pharmaceutical Industry) — fokus pada bagian/section "
                "apa saja yang ada dan bagaimana kontennya disusun"
                if row.category == "format" else
                "bahan sumber DATA untuk menyusun laporan Business Plan Outlook (Global Economic "
                "Outlook, Indonesia Economic Outlook, Pharmaceutical Industry)"
            )
            system = (
                "Kamu adalah analis riset yang meringkas dokumen sumber menjadi poin-poin "
                "terstruktur untuk dipakai berulang kali sebagai referensi oleh AI lain — bukan "
                "narasi panjang. Ringkasan harus padat, faktual, dan mempertahankan semua "
                "angka/statistik/tanggal penting yang ada di dokumen."
            )
            prompt = (
                f'Dokumen berikut adalah {purpose}: "{row.original_name}".\n\n'
                "Ringkas menjadi Markdown bullet list berisi poin-poin kunci saja — angka, tren, "
                "tanggal, dan fakta penting, tanpa basa-basi pembuka/penutup. Gunakan **bold** "
                "untuk angka/istilah kunci. Maksimal sekitar 20 bullet.\n\n"
                "=== ISI DOKUMEN ===\n"
                f"{text_for_ai}{truncation_note}\n"
                "=== AKHIR DOKUMEN ==="
            )

            ai = AIService()
            brief = await ai.complete(system, prompt, provider=provider, gemini_api_key=gemini_api_key)
            brief = (brief or "").strip()
            if not brief:
                raise ValueError("AI tidak mengembalikan ringkasan (respons kosong)")

            row.brief_text = brief
            row.brief_status = "done"
            row.brief_error = None
            row.converted_at = datetime.utcnow()
        except Exception as e:
            logger.warning("outlook_material_convert_failed", material_id=material_id, error=str(e))
            row.brief_status = "failed"
            row.brief_error = str(e)[:2000]

        await db.flush()
        await db.refresh(row)
        return {"success": row.brief_status == "done", "data": self._to_dict(row)}

    def _to_dict(self, row: OutlookMaterial) -> dict:
        return {
            "id":            row.id,
            "plan_year":     row.plan_year,
            "category":      row.category,
            "original_name": row.original_name,
            "content_type":  row.content_type,
            "file_size":     row.file_size,
            "uploaded_by":   row.uploaded_by,
            "created_at":    row.created_at.isoformat() if row.created_at else None,
            "brief_status":  row.brief_status,
            "brief_text":    row.brief_text,
            "brief_error":   row.brief_error,
            "converted_at":  row.converted_at.isoformat() if row.converted_at else None,
        }
