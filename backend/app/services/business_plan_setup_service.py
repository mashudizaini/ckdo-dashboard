"""
Business Plan Setup Service — PAC module.
Stores and retrieves Business Plan Setup documents (Schedule, Guideline, Outlook) in PostgreSQL.
"""
import io
import re
from datetime import datetime, date, timedelta
from typing import Optional
from sqlalchemy import select, delete
from sqlalchemy.ext.asyncio import AsyncSession
from fastapi.responses import StreamingResponse
from app.models.business_plan_setup import PACBusinessPlanSetup
import structlog

logger = structlog.get_logger()

SCHEDULE_DEPTS = [
    ("sales",       "Sales & Marketing"),
    ("development", "Strategic Development"),
    ("plant",        "Plant"),
    ("admin",        "Admin"),
    ("director",     "P. Director"),
]


def _add_markdown_runs(paragraph, text, base_color, accent_color, size=13, italic=False):
    """Split text on **bold** markers into runs, coloring bold segments with
    accent_color — used to render AI-generated Markdown bullets/paragraphs as
    styled PPTX text without a full Markdown renderer."""
    from pptx.util import Pt
    for part in re.split(r"(\*\*[^*]+\*\*)", text):
        if not part:
            continue
        is_bold = part.startswith("**") and part.endswith("**")
        run = paragraph.add_run()
        run.text = part[2:-2] if is_bold else part
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.italic = italic
        run.font.color.rgb = accent_color if is_bold else base_color


class BusinessPlanSetupService:

    # ── List ─────────────────────────────────────────────────────────────────

    async def list_setup(
        self,
        db: AsyncSession,
        setup_module: Optional[str] = None,
        plan_year: Optional[int] = None,
    ) -> dict:
        q = select(PACBusinessPlanSetup).order_by(
            PACBusinessPlanSetup.plan_year.desc(),
            PACBusinessPlanSetup.setup_module,
        )
        if setup_module:
            q = q.where(PACBusinessPlanSetup.setup_module == setup_module)
        if plan_year:
            q = q.where(PACBusinessPlanSetup.plan_year == plan_year)

        result = await db.execute(q)
        rows = result.scalars().all()
        return {"success": True, "count": len(rows), "data": [self._to_dict(r) for r in rows]}

    # ── Get single ────────────────────────────────────────────────────────────

    async def get_setup(self, db: AsyncSession, setup_id: int) -> dict:
        row = await db.get(PACBusinessPlanSetup, setup_id)
        if not row:
            return {"success": False, "error": "Not found"}
        return {"success": True, "data": self._to_dict(row)}

    # ── Upsert ───────────────────────────────────────────────────────────────

    async def upsert_setup(self, db: AsyncSession, payload: dict, username: str) -> dict:
        """
        Create or update a Setup document.
        Unique key: setup_module + plan_year.
        """
        setup_id = payload.get("id")
        row = None

        if setup_id:
            row = await db.get(PACBusinessPlanSetup, setup_id)

        if not row:
            q = select(PACBusinessPlanSetup).where(
                PACBusinessPlanSetup.setup_module == payload.get("setup_module"),
                PACBusinessPlanSetup.plan_year  == payload.get("plan_year"),
            )
            result = await db.execute(q)
            row = result.scalar_one_or_none()

        if row:
            row.content    = payload.get("content", row.content)
            row.status     = payload.get("status",     row.status)
            row.updated_at = datetime.utcnow()
        else:
            row = PACBusinessPlanSetup(
                setup_module = payload.get("setup_module", "schedule"),
                plan_year    = payload.get("plan_year", datetime.now().year),
                content      = payload.get("content", {}),
                status       = payload.get("status", "draft"),
                created_by   = username,
            )
            db.add(row)

        await db.flush()
        await db.refresh(row)
        return {"success": True, "data": self._to_dict(row)}

    # ── Export Schedule to Excel ────────────────────────────────────────────────

    async def export_schedule_excel(self, db: AsyncSession, plan_year: int):
        """Build the Business Plan Schedule Excel, mirroring the layout of
        Business plan schedule.xlsx: title block, a 2-row grouped header
        (Submission Date From/To, Actual Date From/To, PIC per department),
        one row per activity with consecutive-run merging on the date
        columns, and green fill for department cells marked "O"."""
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter

        q = select(PACBusinessPlanSetup).where(
            PACBusinessPlanSetup.setup_module == "schedule",
            PACBusinessPlanSetup.plan_year == plan_year,
        )
        result = await db.execute(q)
        row = result.scalar_one_or_none()
        activities = (row.content or {}).get("activities", []) if row else []

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "BP Schedule"

        thin = Side(style="thin", color="000000")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)
        bold = Font(bold=True)
        center = Alignment(horizontal="center", vertical="center", wrap_text=True)
        left_wrap = Alignment(horizontal="left", vertical="center", wrap_text=True)
        green_fill = PatternFill("solid", fgColor="C6EFCE")
        red_fill = PatternFill("solid", fgColor="FFC7CE")
        hdr_fill = PatternFill("solid", fgColor="D9D9D9")
        red_font = Font(color="9C0006", bold=True)

        N_DEPT = len(SCHEDULE_DEPTS)
        COL_NO, COL_ACT = 1, 2
        COL_PRIOR = 3
        COL_SUB_FROM, COL_SUB_TO = 4, 5
        COL_ACT_FROM, COL_ACT_TO = 6, 7
        COL_DAY = 8
        COL_PIC_START = 9
        COL_PIC_END = COL_PIC_START + N_DEPT - 1
        COL_REQ = COL_PIC_END + 1
        COL_NOTES = COL_REQ + 1
        LAST_COL = COL_NOTES

        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=LAST_COL)
        c = ws.cell(row=1, column=1, value="PT CKD OTTO Pharmaceuticals")
        c.font = Font(bold=True, size=12)
        c.alignment = Alignment(horizontal="left")

        ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=LAST_COL - 2)
        c = ws.cell(row=2, column=1, value=f"Timeline & Schedule / Business Plan {plan_year}")
        c.font = bold
        c.alignment = Alignment(horizontal="left")
        ws.merge_cells(start_row=2, start_column=LAST_COL - 1, end_row=2, end_column=LAST_COL)
        c = ws.cell(row=2, column=LAST_COL - 1, value=date.today().strftime("%b %d, %Y"))
        c.alignment = Alignment(horizontal="right")

        HR1, HR2 = 4, 5

        def merge_hdr(c1, c2, value):
            if c1 == c2:
                ws.merge_cells(start_row=HR1, start_column=c1, end_row=HR2, end_column=c2)
            else:
                ws.merge_cells(start_row=HR1, start_column=c1, end_row=HR1, end_column=c2)
            cell = ws.cell(row=HR1, column=c1, value=value)
            cell.font = bold
            cell.alignment = center
            cell.fill = hdr_fill

        merge_hdr(COL_NO, COL_NO, "No")
        merge_hdr(COL_ACT, COL_ACT, "Schedule")
        merge_hdr(COL_PRIOR, COL_PRIOR, f"Submission Date of {plan_year - 1} BP (in {plan_year - 2})")
        merge_hdr(COL_SUB_FROM, COL_SUB_TO, f"Submission Date of {plan_year} BP (in {plan_year - 1})")
        merge_hdr(COL_ACT_FROM, COL_ACT_TO, "Actual Submission Date")
        merge_hdr(COL_DAY, COL_DAY, "Day")
        merge_hdr(COL_PIC_START, COL_PIC_END, "PIC")
        merge_hdr(COL_REQ, COL_REQ, "Requirement (Form)")
        merge_hdr(COL_NOTES, COL_NOTES, "Notes")

        ws.cell(row=HR2, column=COL_SUB_FROM, value="From").font = bold
        ws.cell(row=HR2, column=COL_SUB_TO, value="To").font = bold
        ws.cell(row=HR2, column=COL_ACT_FROM, value="From").font = bold
        ws.cell(row=HR2, column=COL_ACT_TO, value="To").font = bold
        for i, (_, label) in enumerate(SCHEDULE_DEPTS):
            ws.cell(row=HR2, column=COL_PIC_START + i, value=label).font = bold

        for r in (HR1, HR2):
            for col in range(1, LAST_COL + 1):
                cell = ws.cell(row=r, column=col)
                cell.border = border
                cell.fill = hdr_fill
                if cell.alignment.horizontal is None:
                    cell.alignment = center

        def fmt_date(iso):
            if not iso:
                return None
            try:
                return datetime.strptime(iso, "%Y-%m-%d").date()
            except (ValueError, TypeError):
                return None

        def actual_range(departments):
            dates = sorted(
                d.get("date") for d in (departments or {}).values()
                if d.get("status") == "O" and d.get("date")
            )
            return (dates[0] if dates else None, dates[-1] if dates else None)

        def working_days_diff(d1, d2):
            """Signed working-day count over the range (d1, d2]."""
            if not d1 or not d2:
                return None
            sign = 1
            start, end = d1, d2
            if d2 < d1:
                sign, start, end = -1, d2, d1
            count = 0
            cur = start + timedelta(days=1)
            while cur <= end:
                if cur.weekday() < 5:
                    count += 1
                cur += timedelta(days=1)
            return count * sign

        r0 = HR2 + 1
        rows_data = []
        prev_act_to = None
        for act in activities:
            actual_from, actual_to = actual_range(act.get("departments"))
            act_from_d, act_to_d = fmt_date(actual_from), fmt_date(actual_to)
            note = working_days_diff(prev_act_to, act_from_d) if prev_act_to and act_from_d else None
            prev_act_to = act_to_d
            rows_data.append({
                "no": act.get("no"),
                "activity": act.get("activity", ""),
                "prior": fmt_date(act.get("prior_date")),
                "sub_from": fmt_date(act.get("submission_from")),
                "sub_to": fmt_date(act.get("submission_to")),
                "sub_to_iso": act.get("submission_to"),
                "act_from": act_from_d,
                "act_to": act_to_d,
                "act_to_iso": actual_to,
                "day": act.get("day", ""),
                "departments": act.get("departments") or {},
                "remarks": act.get("remarks", ""),
                "notes": note,
            })

        # Merge consecutive rows sharing the same value, mirroring the
        # source file's boxed grouping on the date columns.
        def merge_run(col, keyfn):
            i = 0
            while i < len(rows_data):
                j = i
                while j + 1 < len(rows_data) and keyfn(rows_data[j + 1]) == keyfn(rows_data[i]) and keyfn(rows_data[i]) is not None:
                    j += 1
                if j > i:
                    ws.merge_cells(start_row=r0 + i, start_column=col, end_row=r0 + j, end_column=col)
                i = j + 1

        merge_run(COL_PRIOR, lambda r: r["prior"])

        for r_idx, rd in enumerate(rows_data, start=r0):
            ws.cell(row=r_idx, column=COL_NO, value=rd["no"])
            ws.cell(row=r_idx, column=COL_ACT, value=rd["activity"])
            ws.cell(row=r_idx, column=COL_PRIOR, value=rd["prior"])
            ws.cell(row=r_idx, column=COL_SUB_FROM, value=rd["sub_from"])
            ws.cell(row=r_idx, column=COL_SUB_TO, value=rd["sub_to"])
            ws.cell(row=r_idx, column=COL_ACT_FROM, value=rd["act_from"])
            ws.cell(row=r_idx, column=COL_ACT_TO, value=rd["act_to"])
            ws.cell(row=r_idx, column=COL_DAY, value=rd["day"])
            ws.cell(row=r_idx, column=COL_REQ, value=rd["remarks"])
            ws.cell(row=r_idx, column=COL_NOTES, value=rd["notes"])

            for col in (COL_PRIOR, COL_SUB_FROM, COL_SUB_TO, COL_ACT_FROM, COL_ACT_TO):
                ws.cell(row=r_idx, column=col).number_format = "dd-mmm-yy"

            act_to_cell = ws.cell(row=r_idx, column=COL_ACT_TO)
            if rd["act_to_iso"] and rd["sub_to_iso"] and rd["act_to_iso"] > rd["sub_to_iso"]:
                act_to_cell.font = red_font

            for i, (key, _) in enumerate(SCHEDULE_DEPTS):
                dept = rd["departments"].get(key, {})
                col = COL_PIC_START + i
                cell = ws.cell(row=r_idx, column=col)
                dept_date_iso = dept.get("date")
                is_late = dept.get("status") == "O" and dept_date_iso and rd["sub_to_iso"] and dept_date_iso > rd["sub_to_iso"]
                if dept.get("status") == "O":
                    if dept_date_iso:
                        cell.value = fmt_date(dept_date_iso)
                        cell.number_format = "dd-mmm-yy"
                    else:
                        cell.value = "O"
                    cell.fill = red_fill if is_late else green_fill
                else:
                    cell.value = "X"

            for col in range(1, LAST_COL + 1):
                cell = ws.cell(row=r_idx, column=col)
                cell.border = border
                cell.alignment = left_wrap if col == COL_ACT else center

        widths = {COL_NO: 5, COL_ACT: 42, COL_PRIOR: 13, COL_SUB_FROM: 12, COL_SUB_TO: 12,
                  COL_ACT_FROM: 12, COL_ACT_TO: 12, COL_DAY: 12, COL_REQ: 14, COL_NOTES: 9}
        for i in range(N_DEPT):
            widths[COL_PIC_START + i] = 12
        for col, w in widths.items():
            ws.column_dimensions[get_column_letter(col)].width = w

        ws.page_setup.orientation = "landscape"
        ws.page_setup.fitToWidth = 1
        ws.sheet_properties.pageSetUpPr.fitToPage = True
        ws.freeze_panes = ws.cell(row=r0, column=COL_ACT + 1)

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        filename = f"Business plan schedule {plan_year}.xlsx"
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    # ── Export Guideline to PowerPoint ──────────────────────────────────────────

    async def export_guideline_ppt(self, db: AsyncSession, plan_year: int):
        """Build the Business Plan Guideline as a PPTX, mirroring Business
        plan guideline.xlsx: a title slide, then one table slide per
        section with Current Year ({plan_year}) / Previous Year
        ({plan_year - 1}) columns."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        q = select(PACBusinessPlanSetup).where(
            PACBusinessPlanSetup.setup_module == "guideline",
            PACBusinessPlanSetup.plan_year == plan_year,
        )
        result = await db.execute(q)
        row = result.scalar_one_or_none()
        sections = (row.content or {}).get("sections", []) if row else []

        NAVY = RGBColor(0x1F, 0x2A, 0x44)
        TEAL = RGBColor(0x0D, 0x94, 0x88)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        HDR_GREY = RGBColor(0xD9, 0xD9, 0xD9)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # ── Title slide ──────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        bg.shadow.inherit = False

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2))
        tf = title_box.text_frame
        tf.text = "Business Plan Guideline"
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8))
        tf = sub_box.text_frame
        tf.text = f"{plan_year} vs {plan_year - 1}"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = TEAL

        company_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.6))
        tf = company_box.text_frame
        tf.text = "PT CKD OTTO Pharmaceuticals"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        # ── One table slide per section ─────────────────────────────────
        for section in sections:
            slide = prs.slides.add_slide(blank)

            head = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
            tf = head.text_frame
            tf.text = f"{section.get('icon', '')}  {section.get('title', '')}".strip()
            tf.paragraphs[0].font.size = Pt(26)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = NAVY

            items = section.get("items", [])
            n_rows = len(items) + 1
            table_shape = slide.shapes.add_table(n_rows, 3, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.7))
            table = table_shape.table
            table.columns[0].width = Inches(6.3)
            table.columns[1].width = Inches(3.0)
            table.columns[2].width = Inches(3.0)

            headers = ["", f"Current Year ({plan_year})", f"Previous Year ({plan_year - 1})"]
            for c, htext in enumerate(headers):
                cell = table.cell(0, c)
                cell.text = htext
                cell.fill.solid()
                cell.fill.fore_color.rgb = HDR_GREY
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(14)
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT

            for r, item in enumerate(items, start=1):
                table.cell(r, 0).text = str(item.get("label", ""))
                table.cell(r, 1).text = str(item.get("current", ""))
                table.cell(r, 2).text = str(item.get("previous", ""))
                for c in range(3):
                    p = table.cell(r, c).text_frame.paragraphs[0]
                    p.font.size = Pt(13)
                    if c:
                        p.alignment = PP_ALIGN.CENTER

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        filename = f"Business plan guideline {plan_year}.pptx"
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    async def export_outlook_ppt(self, db: AsyncSession, plan_year: int):
        """Build the Business Plan Outlook as a management-report-styled
        PPTX: a navy title slide, a Contents slide, then one slide per
        section (Global Economic / Indonesia Economic / Pharmaceutical
        Industry) with a colored header band matching the web view's
        section colors, KPI stat cards pulled from the AI-generated
        **bold** figures, and a two-part body (framing paragraph(s) then
        bullet list) instead of a flat bullet dump."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        q = select(PACBusinessPlanSetup).where(
            PACBusinessPlanSetup.setup_module == "outlook",
            PACBusinessPlanSetup.plan_year == plan_year,
        )
        result = await db.execute(q)
        row = result.scalar_one_or_none()
        content = (row.content or {}) if row else {}

        NAVY = RGBColor(0x1F, 0x2A, 0x44)
        TEAL = RGBColor(0x0D, 0x94, 0x88)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
        CARD_BORDER = RGBColor(0xE2, 0xE5, 0xEA)
        TEXT_DARK = RGBColor(0x1F, 0x2A, 0x44)
        TEXT_GRAY = RGBColor(0x6B, 0x72, 0x80)

        SECTIONS = [
            {"key": "global_economic", "icon": "🌍", "accent": RGBColor(0x25, 0x63, 0xEB), "accent_light": RGBColor(0xE8, 0xEF, 0xFD)},
            {"key": "indonesia_economic", "icon": "🏛️", "accent": RGBColor(0xD9, 0x77, 0x06), "accent_light": RGBColor(0xFC, 0xF1, 0xE1)},
            {"key": "pharmaceutical", "icon": "💊", "accent": RGBColor(0x05, 0x96, 0x69), "accent_light": RGBColor(0xE1, 0xF5, 0xEE)},
        ]

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        SW, SH = prs.slide_width, prs.slide_height

        def add_footer(slide, page_num):
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.08), Inches(12.33), Pt(1))
            line.fill.solid(); line.fill.fore_color.rgb = CARD_BORDER; line.line.fill.background(); line.shadow.inherit = False
            f = slide.shapes.add_textbox(Inches(0.5), Inches(7.13), Inches(9), Inches(0.3))
            p = f.text_frame.paragraphs[0]
            p.text = "PT CKD OTTO Pharmaceuticals — Business Plan Outlook · Internal / Management Use"
            p.font.size = Pt(9); p.font.color.rgb = TEXT_GRAY
            f2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.13), Inches(1.3), Inches(0.3))
            p2 = f2.text_frame.paragraphs[0]
            p2.text = str(page_num)
            p2.font.size = Pt(9); p2.font.color.rgb = TEXT_GRAY; p2.alignment = PP_ALIGN.RIGHT

        # ── Title slide ──────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False

        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.75), Inches(2.2), Pt(4))
        accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = TEAL; accent_bar.line.fill.background(); accent_bar.shadow.inherit = False

        company_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.6))
        tf = company_box.text_frame
        tf.text = "PT CKD OTTO Pharmaceuticals"
        tf.paragraphs[0].font.size = Pt(16); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.5))
        tf = title_box.text_frame; tf.word_wrap = True
        tf.text = "Business Plan Outlook"
        tf.paragraphs[0].font.size = Pt(46); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE

        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.85), Inches(11), Inches(0.7))
        tf = sub_box.text_frame
        tf.text = f"Economic & Industry Outlook — {plan_year}"
        tf.paragraphs[0].font.size = Pt(20); tf.paragraphs[0].font.color.rgb = RGBColor(0x9C, 0xC5, 0xFF)

        meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
        tf = meta_box.text_frame
        tf.text = "Prepared for Management Review · Internal Use Only"
        tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.italic = True; tf.paragraphs[0].font.color.rgb = RGBColor(0x8C, 0x96, 0xB0)

        # ── Contents slide ───────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_BG; bg.line.fill.background(); bg.shadow.inherit = False

        head = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.8))
        tf = head.text_frame
        tf.text = "Contents"
        tf.paragraphs[0].font.size = Pt(32); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = NAVY

        y = 2.0
        for sec in SECTIONS:
            section = content.get(sec["key"]) or {}
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.7), Inches(1.2))
            card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = CARD_BORDER; card.line.width = Pt(0.75); card.shadow.inherit = False
            try:
                card.adjustments[0] = 0.08
            except Exception:
                pass

            accent_dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Pt(5), Inches(1.2))
            accent_dot.fill.solid(); accent_dot.fill.fore_color.rgb = sec["accent"]; accent_dot.line.fill.background(); accent_dot.shadow.inherit = False

            icon_box = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.28), Inches(0.8), Inches(0.6))
            icon_box.text_frame.text = sec["icon"]
            icon_box.text_frame.paragraphs[0].font.size = Pt(28)

            title_box = slide.shapes.add_textbox(Inches(2.1), Inches(y + 0.35), Inches(9.8), Inches(0.6))
            tf = title_box.text_frame; tf.word_wrap = True
            tf.text = section.get("title", sec["key"])
            tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = sec["accent"]

            y += 1.5

        add_footer(slide, 2)

        # ── One slide per section ───────────────────────────────────────
        for page_idx, sec in enumerate(SECTIONS, start=3):
            section = content.get(sec["key"]) or {}
            accent = sec["accent"]
            accent_light = sec["accent_light"]
            slide = prs.slides.add_slide(blank)
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
            bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background(); bg.shadow.inherit = False

            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.05))
            band.fill.solid(); band.fill.fore_color.rgb = accent; band.line.fill.background(); band.shadow.inherit = False

            icon_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(0.7), Inches(0.7))
            icon_box.text_frame.text = sec["icon"]
            icon_box.text_frame.paragraphs[0].font.size = Pt(30)

            head = slide.shapes.add_textbox(Inches(1.25), Inches(0.22), Inches(11.3), Inches(0.65))
            tf = head.text_frame; tf.word_wrap = True
            tf.text = section.get("title", "")
            tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE

            raw_text = (section.get("text") or "").strip()
            lines = [ln.strip() for ln in raw_text.split("\n") if ln.strip()]
            bullet_lines = [ln for ln in lines if ln.startswith("- ") or ln.startswith("* ")]
            para_lines = [ln for ln in lines if ln not in bullet_lines]

            # KPI stat cards — take the first bullets that carry two **bold**
            # groups (label, value); this is the only visual we derive from
            # the freeform text, so it's limited to figures the AI already
            # committed to in writing rather than any fabricated chart.
            stats = []
            for ln in bullet_lines:
                bolds = re.findall(r"\*\*([^*]+)\*\*", ln)
                if len(bolds) >= 2 and len(bolds[1]) <= 24:
                    stats.append((bolds[0], bolds[1]))
                if len(stats) >= 3:
                    break

            content_top = 1.3
            if stats:
                card_w, gap, card_h = Inches(3.75), Inches(0.2), Inches(1.0)
                x = Inches(0.5)
                for label, value in stats:
                    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(content_top), card_w, card_h)
                    card.fill.solid(); card.fill.fore_color.rgb = accent_light; card.line.fill.background(); card.shadow.inherit = False
                    try:
                        card.adjustments[0] = 0.12
                    except Exception:
                        pass
                    tf = card.text_frame; tf.word_wrap = True
                    tf.margin_left = Pt(12); tf.margin_right = Pt(12); tf.margin_top = Pt(8); tf.margin_bottom = Pt(6)
                    p0 = tf.paragraphs[0]
                    p0.text = value
                    p0.font.size = Pt(20); p0.font.bold = True; p0.font.color.rgb = accent
                    p1 = tf.add_paragraph()
                    p1.text = label
                    p1.font.size = Pt(10.5); p1.font.color.rgb = TEXT_GRAY
                    x = x + card_w + gap
                content_top = 2.5

            body_box = slide.shapes.add_textbox(Inches(0.6), Inches(content_top), Inches(12.1), Inches(7.0 - content_top))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True

            first_para = True
            for ln in para_lines:
                p = body_tf.paragraphs[0] if first_para else body_tf.add_paragraph()
                first_para = False
                p.space_after = Pt(10)
                _add_markdown_runs(p, ln, TEXT_DARK, accent, size=13, italic=True)

            for ln in bullet_lines:
                p = body_tf.paragraphs[0] if first_para else body_tf.add_paragraph()
                first_para = False
                p.space_after = Pt(8)
                bullet_text = re.sub(r"^[-*]\s+", "", ln)
                bullet_run = p.add_run()
                bullet_run.text = "●  "
                bullet_run.font.size = Pt(13)
                bullet_run.font.color.rgb = accent
                _add_markdown_runs(p, bullet_text, TEXT_DARK, accent, size=13)

            add_footer(slide, page_idx)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        filename = f"Business plan outlook {plan_year}.pptx"
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )

    # ── Delete ────────────────────────────────────────────────────────────────

    async def delete_setup(self, db: AsyncSession, setup_id: int) -> dict:
        row = await db.get(PACBusinessPlanSetup, setup_id)
        if not row:
            return {"success": False, "error": "Not found"}
        await db.execute(delete(PACBusinessPlanSetup).where(PACBusinessPlanSetup.id == setup_id))
        return {"success": True, "message": f"Deleted setup #{setup_id}"}

    # ── Helper ────────────────────────────────────────────────────────────────

    def _to_dict(self, row: PACBusinessPlanSetup) -> dict:
        return {
            "id":          row.id,
            "setup_module": row.setup_module,
            "plan_year":   row.plan_year,
            "content":     row.content,
            "status":      row.status,
            "created_at":  row.created_at.isoformat() if row.created_at else None,
            "updated_at":  row.updated_at.isoformat() if row.updated_at else None,
            "created_by":  row.created_by,
        }
